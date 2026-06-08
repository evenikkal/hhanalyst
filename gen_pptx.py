"""Generate a 6-slide, ~2-minute defense presentation for the hhanalyst course paper."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = "E:/mitp/hhanalyst"
DIAG = f"{ROOT}/diagrams"
OUT = f"{ROOT}/Защита_hhanalyst.pptx"

# Palette (matches the paper's diagram colors)
BLUE = RGBColor(0x2E, 0x5F, 0xA3)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE6, 0x51, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Calibri"


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def txt(s, x, y, w, h, lines, size=18, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=6, font=FONT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(ln, tuple):
            text, lvl, lsize, lbold, lcolor = ln
        else:
            text, lvl, lsize, lbold, lcolor = ln, 0, size, bold, color
        p.level = lvl
        r = p.add_run()
        r.text = text
        r.font.size = Pt(lsize)
        r.font.bold = lbold
        r.font.color.rgb = lcolor
        r.font.name = font
    return tb


def bullets(s, x, y, w, h, items, size=18, color=DARK, gap=10):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        r = p.add_run()
        r.text = ("•  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.name = FONT
    return tb


def header(s, num, title):
    """Accent bar + title for content slides."""
    rect(s, 0, 0, SW, Inches(1.15), BLUE)
    rect(s, 0, Inches(1.15), SW, Pt(4), ORANGE)
    txt(s, Inches(0.6), Inches(0.12), Inches(11), Inches(0.9),
        title, size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # slide number badge
    txt(s, SW - Inches(1.0), Inches(0.12), Inches(0.7), Inches(0.9),
        f"{num}/6", size=16, color=RGBColor(0xCF, 0xDD, 0xF2), bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def pic_fit(s, path, x, y, max_w, max_h):
    """Add picture scaled to fit box, centered."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    s.shapes.add_picture(path, px, py, width=w, height=h)


# ---------------------------------------------------------------- Slide 1: Title
s = slide()
rect(s, 0, 0, SW, SH, BLUE)
rect(s, 0, Inches(4.55), SW, Pt(4), ORANGE)
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.5),
    "Защита курсовой работы", size=18, color=RGBColor(0xCF, 0xDD, 0xF2),
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.4),
    ["Веб-сервис анализа рынка IT-вакансий hhanalyst",
     ("сбор, NLP-обработка и визуализация данных hh.ru", 0, 22, False,
      RGBColor(0xDD, 0xE8, 0xF6))],
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE, space_after=14)
txt(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(2.0),
    ["Дисциплина: «Технологии разработки программного обеспечения»",
     "Выполнил: студент группы ________",
     "Москва, 2026"],
    size=18, color=RGBColor(0xDD, 0xE8, 0xF6), align=PP_ALIGN.CENTER, space_after=8)

# ---------------------------------------------------------------- Slide 2: Problem & goal
s = slide()
header(s, 2, "Проблема и цель")
bullets(s, Inches(0.6), Inches(1.5), Inches(7.4), Inches(5.5), [
    "Анализ требований IT-рынка вручную — медленно и неактуально",
    "HR-специалисту и соискателю нужны живые данные: какие навыки и уровни востребованы",
    ("Цель: разработать веб-сервис сбора и NLP-анализа вакансий с hh.ru", 0),
    ("собрать вакансии по запросу через API hh.ru", 1),
    ("извлечь навыки средствами NLP для русского языка", 1),
    ("классифицировать уровень: intern / junior / middle / senior", 1),
    ("построить наглядную визуализацию результатов", 1),
], size=20, gap=12)
# side accent panel
rect(s, Inches(8.5), Inches(1.6), Inches(4.2), Inches(5.0), LIGHT)
txt(s, Inches(8.8), Inches(1.9), Inches(3.6), Inches(4.4),
    [("Входные данные", 0, 16, True, BLUE),
     ("вакансии hh.ru по поисковому запросу", 0, 15, False, GREY),
     ("", 0, 6, False, GREY),
     ("Результат", 0, 16, True, BLUE),
     ("топ востребованных навыков", 0, 15, False, GREY),
     ("распределение по уровням", 0, 15, False, GREY),
     ("тепловая карта «навык × уровень»", 0, 15, False, GREY)],
    space_after=6)

# ---------------------------------------------------------------- Slide 3: Architecture
s = slide()
header(s, 3, "Архитектура системы")
pic_fit(s, f"{DIAG}/arch.png", Inches(0.4), Inches(1.35), Inches(8.3), Inches(5.7))
rect(s, Inches(8.9), Inches(1.5), Inches(4.0), Inches(5.3), LIGHT)
txt(s, Inches(9.15), Inches(1.7), Inches(3.5), Inches(5.0),
    [("Два микросервиса", 0, 18, True, BLUE),
     ("Go-сборщик", 0, 16, True, DARK),
     ("пул из 5 воркеров, rate-limiter 4 req/s, образ Docker ~8 МБ", 0, 14, False, GREY),
     ("", 0, 4, False, GREY),
     ("Python-аналитик", 0, 16, True, DARK),
     ("FastAPI + Natasha (NLP)", 0, 14, False, GREY),
     ("", 0, 6, False, GREY),
     ("Почему два языка?", 0, 16, True, ORANGE),
     ("Go — конкурентный сетевой сбор; Python — зрелые NLP-библиотеки", 0, 14, False, GREY)],
    space_after=6)

# ---------------------------------------------------------------- Slide 4: Key tech decision
s = slide()
header(s, 4, "Ключевое техническое решение")
bullets(s, Inches(0.6), Inches(1.5), Inches(7.3), Inches(5.3), [
    "Двухуровневый кэш: L1 в памяти + L2 на SQLite",
    "6-уровневая цепочка fallback → сервис работает даже офлайн",
    "Кэш экономит дорогой NLP-проход при повторных запросах",
    ("свежие данные отдаются мгновенно из L1/L2", 1),
    ("устаревшие — фоново обновляются, ответ не блокируется", 1),
], size=20, gap=14)
# metric cards
def card(x, big, small, c):
    rect(s, x, Inches(4.7), Inches(3.7), Inches(2.1), LIGHT)
    txt(s, x, Inches(4.85), Inches(3.7), Inches(1.0), big, size=40, color=c,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, Inches(5.95), Inches(3.7), Inches(0.8), small, size=15, color=GREY,
        align=PP_ALIGN.CENTER)
rect(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(2.9), LIGHT)
txt(s, Inches(8.55), Inches(1.7), Inches(4.0), Inches(2.6),
    [("Эффект кэширования", 0, 18, True, BLUE),
     ("Чтение из L2 (SQLite):", 0, 15, False, DARK),
     ("0.38 мс", 0, 30, True, GREEN),
     ("Полный NLP-проход:", 0, 15, False, DARK),
     ("≈ 670 мс", 0, 30, True, ORANGE)],
    space_after=4)
txt(s, Inches(8.3), Inches(4.7), Inches(4.5), Inches(2.1),
    [("Ускорение повторного запроса", 0, 16, True, BLUE),
     ("≈ в 1700 раз", 0, 34, True, GREEN)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=6)

# ---------------------------------------------------------------- Slide 5: Results
s = slide()
header(s, 5, "Результаты анализа")
pic_fit(s, f"{DIAG}/chart_top_skills.png", Inches(0.3), Inches(1.4), Inches(6.5), Inches(4.6))
pic_fit(s, f"{DIAG}/chart_levels.png", Inches(6.9), Inches(1.4), Inches(6.1), Inches(4.6))
rect(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(1.0), LIGHT)
txt(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.8),
    [("Производительность NLP-конвейера: ≈ 74 мс на одну вакансию  •  "
      "извлечение навыков 11 мс  •  классификация уровня 1.4 мс", 0, 16, True, DARK)],
    anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- Slide 6: Conclusion
s = slide()
header(s, 6, "Итоги работы")
bullets(s, Inches(0.6), Inches(1.5), Inches(7.4), Inches(5.0), [
    "Цель достигнута: работающий веб-сервис анализа IT-вакансий",
    "Два микросервиса на Go и Python, общение по REST API",
    "Автоматические тесты (Go + Python) и CI/CD на GitHub Actions",
    "Контейнеризация: Docker-образ сборщика ~8 МБ",
    "Надёжность за счёт кэша и fallback — работа даже офлайн",
], size=20, gap=14)
rect(s, Inches(8.5), Inches(1.6), Inches(4.3), Inches(4.6), BLUE)
txt(s, Inches(8.8), Inches(2.6), Inches(3.7), Inches(2.6),
    ["Спасибо за внимание!",
     ("Готов ответить на вопросы", 0, 18, False, RGBColor(0xDD, 0xE8, 0xF6))],
    size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE, space_after=12)

# ---------------------------------------------------------------- Speaker notes (~2 min)
notes = [
    # 1
    "Здравствуйте. Тема моей курсовой работы — веб-сервис анализа рынка "
    "IT-вакансий hhanalyst: сбор, NLP-обработка и визуализация данных с hh.ru.",
    # 2
    "Анализировать требования IT-рынка вручную долго и быстро устаревает. "
    "И HR-специалисту, и соискателю нужны живые данные: какие навыки и какие "
    "уровни сейчас востребованы. Поэтому цель работы — разработать сервис, "
    "который сам собирает вакансии с hh.ru, средствами NLP извлекает навыки, "
    "классифицирует уровень от intern до senior и наглядно визуализирует результат.",
    # 3
    "Система состоит из двух микросервисов. Сборщик на Go использует пул из пяти "
    "воркеров и rate-limiter на четыре запроса в секунду — это быстрый "
    "конкурентный сбор данных. Аналитик на Python с FastAPI и библиотекой "
    "Natasha отвечает за NLP. Два языка выбраны осознанно: Go силён в "
    "конкурентном сетевом сборе, а у Python зрелые NLP-библиотеки для русского языка.",
    # 4
    "Ключевое решение — двухуровневый кэш: в памяти и на SQLite, плюс "
    "шестиуровневая цепочка fallback, благодаря которой сервис работает даже "
    "офлайн. Это важно, потому что полный NLP-проход стоит около 670 "
    "миллисекунд, а чтение из кэша — всего 0,38 миллисекунды. То есть повторный "
    "запрос обрабатывается почти в полторы тысячи раз быстрее.",
    # 5
    "Вот результаты на реальной выборке: слева — топ востребованных навыков, "
    "справа — распределение вакансий по уровням. По производительности: "
    "конвейер обрабатывает одну вакансию примерно за 74 миллисекунды, "
    "извлечение навыков — 11 миллисекунд, классификация уровня — около полутора.",
    # 6
    "Подведу итог. Цель достигнута: получился рабочий сервис из двух "
    "микросервисов на Go и Python, со связью по REST API, автотестами, "
    "настроенным CI/CD на GitHub Actions и контейнеризацией — образ сборщика "
    "около 8 мегабайт. Спасибо за внимание, готов ответить на вопросы.",
]
for sl, note in zip(prs.slides, notes):
    sl.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
